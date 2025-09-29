#!/usr/bin/env python3
"""
CLI tool to embed documents (Markdown, JSON, PowerPoint, Word, PDF, Text) using OpenAI embeddings
and store them in categorized Pinecone indices.

Usage:
    embedder embed docs/
    embedder search "What day is the cold plunge?" --index idx-ops
    embedder stats --index idx-program
    embedder clear --namespace default --index idx-medical

Requirements:
- pip install openai pinecone-client python-dotenv click python-pptx python-docx PyPDF2 tiktoken
- Set OPENAI_API_KEY and PINECONE_API_KEY in environment or .env file

Indices:
- idx-ops: Operations (logistics, pricing, venue, booking)
- idx-legal: Legal & compliance (privacy, licensing, terms)
- idx-medical: Health & safety (screening, medications, contraindications)
- idx-program: Program content (retreats, ceremonies, modalities, schedule)
"""

import os
import glob
import json
import hashlib
import sys
import re
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from datetime import datetime
import click
from dotenv import load_dotenv
from openai import OpenAI
from pinecone import Pinecone
import tiktoken

# Document processing libraries
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

# Load environment variables
load_dotenv()

class DocumentEmbedder:
    def __init__(self):
        # Initialize OpenAI client
        self.openai_client = OpenAI(
            api_key=os.getenv('OPENAI_API_KEY')
        )

        # Initialize Pinecone
        self.pc = Pinecone(
            api_key=os.getenv('PINECONE_API_KEY')
        )

        # Configuration
        self.embedding_model = "text-embedding-3-small"
        self.embedding_dimension = 1536
        self.default_namespace = "default"

        # Default configuration - can be overridden by user input
        self.index_name = None  # Will be set by user input

        # Default category-based namespaces
        self.default_namespaces = {
            "ops": "Operations (logistics, pricing, venue, booking)",
            "legal": "Legal & compliance (privacy, licensing, terms)",
            "medical": "Health & safety (screening, medications, contraindications)",
            "program": "Program content (retreats, ceremonies, modalities, schedule)"
        }

        # Will be set by user selection
        self.selected_namespaces = {}

        # Active connections to indices
        self.index_connections = {}

        # Tokenizer for chunk sizing
        self.tokenizer = tiktoken.get_encoding("cl100k_base")

        # Categorization keywords
        self.category_keywords = {
            "ops": {
                "paths": ["location", "venue", "contact", "home", "start-application", "book-a-journey", "join", "journeys", "pricing", "deposit"],
                "keywords": ["address", "check-in", "deposit", "refund", "guarantee", "room", "amenities", "accessibility", "parking", "transportation", "apply", "book", "dates", "schedule", "arrival", "venue", "location", "logistics", "what to bring"]
            },
            "legal": {
                "paths": ["privacy-policy", "training-licensed-facilitator-colorado", "terms", "consent", "licensing"],
                "keywords": ["privacy", "consent", "licensing", "colorado", "compliance", "legal", "terms", "conditions", "liability", "agreement", "data", "retention", "rfra", "religious freedom"]
            },
            "medical": {
                "paths": ["resources-health-intake", "health", "screening", "medical"],
                "keywords": ["screening", "health", "medical", "medication", "ssri", "maoi", "contraindication", "pregnant", "heart condition", "emergency", "safety", "intake", "eligibility", "doctor", "physician", "health condition"]
            },
            "program": {
                "paths": ["journeys-awaken", "journeys-open", "journeys-visionary", "journeys-program", "about-values", "about-beliefs", "about-sacred-container", "about-sacraments", "community", "live", "courses", "training"],
                "keywords": ["ceremony", "ceremonies", "retreat", "program", "schedule", "daily", "breathwork", "yoga", "integration", "modalities", "ayahuasca", "mushrooms", "psilocybin", "sacrament", "spiritual", "meditation", "values", "beliefs", "sacred container"]
            }
        }

    def connect_to_index(self, index_name: str = None):
        """Connect to a specific Pinecone index."""
        target_index = index_name or self.index_name
        if not target_index:
            raise ValueError("No index name specified")

        if target_index not in self.index_connections:
            self.index_connections[target_index] = self.pc.Index(target_index)
        return self.index_connections[target_index]

    def get_available_indices(self) -> List[str]:
        """Get list of available Pinecone indices."""
        try:
            indexes = self.pc.list_indexes()
            return [idx.name for idx in indexes]
        except Exception as e:
            print(f"Error fetching indices: {e}")
            return []

    def prompt_for_index_selection(self) -> str:
        """Prompt user to select a Pinecone index."""
        available_indices = self.get_available_indices()

        if not available_indices:
            raise Exception("No Pinecone indices found")

        print("\n📊 Available Pinecone indices:")
        for i, idx_name in enumerate(available_indices, 1):
            default_marker = " (default)" if idx_name == "journeys-faq" else ""
            print(f"  {i}. {idx_name}{default_marker}")

        while True:
            response = input(f"\nSelect index (1-{len(available_indices)}) or press Enter for 'journeys-faq': ").strip()

            if not response and "journeys-faq" in available_indices:
                return "journeys-faq"

            try:
                idx = int(response) - 1
                if 0 <= idx < len(available_indices):
                    return available_indices[idx]
                else:
                    print(f"Please enter a number between 1 and {len(available_indices)}")
            except ValueError:
                if response in available_indices:
                    return response
                print("Please enter a valid number or index name")

    def prompt_for_namespace_selection(self) -> Dict[str, str]:
        """Prompt user to select which namespaces to use."""
        print("\n📁 Default namespaces for categorization:")
        for key, desc in self.default_namespaces.items():
            print(f"  • {key}: {desc}")

        print("\nOptions:")
        print("  1. Use all default namespaces (recommended)")
        print("  2. Select specific namespaces")
        print("  3. Use custom namespaces")

        while True:
            choice = input("\nSelect option (1-3) or press Enter for default: ").strip()

            if not choice or choice == "1":
                return self.default_namespaces.copy()

            elif choice == "2":
                selected = {}
                print("\nSelect namespaces to include:")
                for key, desc in self.default_namespaces.items():
                    include = input(f"Include '{key}' ({desc})? [Y/n]: ").strip().lower()
                    if include in ['', 'y', 'yes']:
                        selected[key] = desc
                return selected if selected else self.default_namespaces.copy()

            elif choice == "3":
                custom = {}
                print("\nEnter custom namespaces (press Enter with empty name to finish):")
                while True:
                    name = input("Namespace name: ").strip()
                    if not name:
                        break
                    desc = input(f"Description for '{name}': ").strip()
                    custom[name] = desc or f"Custom namespace: {name}"
                return custom if custom else self.default_namespaces.copy()

            else:
                print("Please enter 1, 2, or 3")

    def read_file_content(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Read content from various file types and return content with metadata."""
        file_ext = Path(file_path).suffix.lower()
        metadata = {"source_file": file_path, "file_type": file_ext}

        try:
            if file_ext == '.md':
                return self._read_markdown(file_path), metadata
            elif file_ext == '.json':
                return self._read_json(file_path, metadata)
            elif file_ext in ['.ppt', '.pptx'] and PPTX_AVAILABLE:
                return self._read_powerpoint(file_path), metadata
            elif file_ext in ['.doc', '.docx'] and DOCX_AVAILABLE:
                return self._read_word(file_path), metadata
            elif file_ext == '.pdf' and PDF_AVAILABLE:
                return self._read_pdf(file_path), metadata
            elif file_ext == '.txt':
                return self._read_text(file_path), metadata
            else:
                print(f"Unsupported file type: {file_ext} for {file_path}")
                return "", metadata
        except Exception as e:
            print(f"Error reading {file_path}: {e}")
            return "", metadata

    def _read_markdown(self, file_path: str) -> str:
        """Read markdown file content."""
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()

    def _read_text(self, file_path: str) -> str:
        """Read text file content."""
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()

    def _read_json(self, file_path: str, metadata: Dict[str, Any]) -> Tuple[str, Dict[str, Any]]:
        """Read and process JSON website content."""
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)

        # Extract page metadata
        if 'page' in data:
            page = data['page']
            metadata.update({
                'page_slug': page.get('slug', ''),
                'page_title': page.get('title', ''),
                'page_route': page.get('route', ''),
                'source_url': f"https://ceremoniacircle.org{page.get('route', '')}"
            })

            # Process sections and blocks
            content_parts = []
            if page.get('title'):
                content_parts.append(f"# {page['title']}")

            for section in page.get('sections', []):
                section_content = self._extract_section_content(section)
                if section_content:
                    content_parts.append(section_content)

            return '\n\n'.join(content_parts), metadata

        return json.dumps(data, indent=2), metadata

    def _extract_section_content(self, section: Dict[str, Any]) -> str:
        """Extract readable content from a JSON section."""
        content_parts = []

        for block in section.get('blocks', []):
            block_type = block.get('type', '')
            props = block.get('props', {})

            if block_type == 'heading' and 'text' in props:
                level = props.get('level', 2)
                heading_prefix = '#' * level
                content_parts.append(f"{heading_prefix} {props['text']}")

            elif block_type == 'text' and 'text' in props:
                content_parts.append(props['text'])

            elif block_type == 'image' and 'alt' in props and props['alt']:
                content_parts.append(f"[Image: {props['alt']}]")

        return '\n\n'.join(content_parts)

    def _read_powerpoint(self, file_path: str) -> str:
        """Read PowerPoint file content."""
        prs = Presentation(file_path)
        content_parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = [f"## Slide {slide_num}"]

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())

            if len(slide_content) > 1:  # Has content beyond slide number
                content_parts.append('\n'.join(slide_content))

        return '\n\n'.join(content_parts)

    def _read_word(self, file_path: str) -> str:
        """Read Word document content."""
        doc = Document(file_path)
        content_parts = []

        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                content_parts.append(paragraph.text.strip())

        return '\n\n'.join(content_parts)

    def _read_pdf(self, file_path: str) -> str:
        """Read PDF file content."""
        content_parts = []

        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)

            for page_num, page in enumerate(pdf_reader.pages, 1):
                try:
                    text = page.extract_text()
                    if text.strip():
                        content_parts.append(f"## Page {page_num}\n{text.strip()}")
                except Exception as e:
                    print(f"Error extracting text from page {page_num} of {file_path}: {e}")

        return '\n\n'.join(content_parts)

    def chunk_text_by_tokens(self, text: str, max_tokens: int = 500, overlap_tokens: int = 80) -> List[str]:
        """
        Split text into overlapping chunks by token count for better embedding coverage.
        """
        tokens = self.tokenizer.encode(text)

        if len(tokens) <= max_tokens:
            return [text]

        chunks = []
        start = 0

        while start < len(tokens):
            end = min(start + max_tokens, len(tokens))

            # Get the chunk tokens
            chunk_tokens = tokens[start:end]
            chunk_text = self.tokenizer.decode(chunk_tokens)

            # Try to break at sentence boundaries if not at the end
            if end < len(tokens):
                # Look for sentence endings in the last portion of the chunk
                sentences = re.split(r'[.!?]\s+', chunk_text)
                if len(sentences) > 1:
                    # Keep all but the last incomplete sentence
                    chunk_text = '. '.join(sentences[:-1]) + '.'

            if chunk_text.strip():
                chunks.append(chunk_text.strip())

            # Move start position with overlap
            if end >= len(tokens):
                break

            start = end - overlap_tokens
            if start >= len(tokens):
                break

        return chunks

    def extract_sections_from_content(self, content: str, file_path: str, file_metadata: Dict[str, Any]) -> List[Dict[str, Any]]:
        """
        Extract meaningful sections from content with enhanced metadata.
        """
        sections = []
        file_name = Path(file_path).stem
        page_slug = file_metadata.get('page_slug', '')

        # Determine category for this content
        category = self.categorize_content(file_path, page_slug, content)

        # Split content into logical sections
        if file_metadata.get('file_type') == '.json':
            # For JSON files, treat as single section but extract heading path
            heading_path = self._extract_heading_path(content, file_metadata)
            sections.append({
                "title": file_metadata.get('page_title', file_name),
                "content": content,
                "category": category,
                "page_slug": page_slug,
                "heading_path": heading_path,
                "source_file": file_path,
                "file_name": file_name,
                "file_metadata": file_metadata
            })
        else:
            # For other files, split by headers
            sections = self._split_by_headers(content, file_path, category, file_metadata)

        return sections

    def _extract_heading_path(self, content: str, file_metadata: Dict[str, Any]) -> str:
        """Extract heading path for structured navigation."""
        page_title = file_metadata.get('page_title', '')
        page_slug = file_metadata.get('page_slug', '')

        # Build path from page structure
        path_parts = []

        # Add top-level category based on slug
        if 'journey' in page_slug:
            path_parts.append('Journeys')
        elif 'about' in page_slug:
            path_parts.append('About')
        elif 'training' in page_slug:
            path_parts.append('Training')
        elif 'course' in page_slug:
            path_parts.append('Courses')

        if page_title:
            path_parts.append(page_title)

        return ' › '.join(path_parts) if path_parts else page_slug

    def _split_by_headers(self, content: str, file_path: str, category: str, file_metadata: Dict[str, Any]) -> List[Dict[str, Any]]:
        """Split content by markdown headers or logical breaks."""
        sections = []
        file_name = Path(file_path).stem
        lines = content.split('\n')
        current_section = {"title": "", "content": "", "level": 0}

        for line in lines:
            line = line.strip()

            # Skip empty lines and bookmarks
            if not line or line.startswith('[bookmark:') or line.startswith('[image:'):
                continue

            # Check for markdown headers or section-like content
            is_header = False
            if line.startswith('#'):
                is_header = True
                header_level = len(line) - len(line.lstrip('#'))
                title = line.lstrip('# ').strip()
            elif (line and not line.startswith(' ') and
                  (line.isupper() or
                   any(keyword in line.lower() for keyword in ['day ', 'logistics', 'team', 'ceremony', 'meditation', 'workshop']) or
                   line.endswith(':') and len(line) < 100)):
                is_header = True
                title = line

            if is_header:
                # Save previous section if it has content
                if current_section["content"].strip():
                    sections.append(self._create_section_dict(
                        current_section["title"] or f"{file_name} Section",
                        current_section["content"].strip(),
                        category, file_path, file_name, file_metadata
                    ))

                # Start new section
                current_section = {"title": title, "content": "", "level": 1}
            else:
                # Add to current section content
                if current_section["content"]:
                    current_section["content"] += "\n"
                current_section["content"] += line

        # Add final section
        if current_section["content"].strip():
            sections.append(self._create_section_dict(
                current_section["title"] or f"{file_name} Section",
                current_section["content"].strip(),
                category, file_path, file_name, file_metadata
            ))

        # If no sections found, treat entire document as one section
        if not sections:
            sections.append(self._create_section_dict(
                f"{file_name} Complete Document",
                content, category, file_path, file_name, file_metadata
            ))

        return sections

    def _create_section_dict(self, title: str, content: str, category: str,
                           file_path: str, file_name: str, file_metadata: Dict[str, Any]) -> Dict[str, Any]:
        """Create a standardized section dictionary with full metadata."""
        return {
            "title": title,
            "content": content,
            "category": category,
            "page_slug": file_metadata.get('page_slug', ''),
            "heading_path": file_metadata.get('page_title', title),
            "source_file": file_path,
            "file_name": file_name,
            "source_url": file_metadata.get('source_url', ''),
            "file_metadata": file_metadata
        }

    def categorize_content(self, file_path: str, page_slug: str = "", content: str = "") -> str:
        """Categorize content into one of the four indices based on path, slug, and content analysis."""
        file_path_lower = file_path.lower()
        page_slug_lower = page_slug.lower()
        content_lower = content.lower()

        # Score each category
        scores = {"ops": 0, "legal": 0, "medical": 0, "program": 0}

        for category, rules in self.category_keywords.items():
            # Check file path patterns
            for path_pattern in rules["paths"]:
                if path_pattern in file_path_lower or path_pattern in page_slug_lower:
                    scores[category] += 3

            # Check content keywords
            for keyword in rules["keywords"]:
                if keyword in content_lower:
                    scores[category] += 1
                # Boost score for keywords in prominent positions
                if keyword in content_lower[:500]:  # First 500 chars
                    scores[category] += 1

        # Default fallback based on common patterns
        if "health" in file_path_lower or "medical" in file_path_lower:
            scores["medical"] += 5
        elif "legal" in file_path_lower or "privacy" in file_path_lower:
            scores["legal"] += 5
        elif "location" in file_path_lower or "booking" in file_path_lower:
            scores["ops"] += 5
        elif "journey" in file_path_lower or "retreat" in file_path_lower or "program" in file_path_lower:
            scores["program"] += 5

        # Return category with highest score, default to program
        best_category = max(scores.items(), key=lambda x: x[1])
        return best_category[0] if best_category[1] > 0 else "program"

    def generate_embedding(self, text: str) -> List[float]:
        """Generate embedding using OpenAI text-embedding-3-small."""
        try:
            response = self.openai_client.embeddings.create(
                model=self.embedding_model,
                input=text,
                encoding_format="float"
            )
            return response.data[0].embedding
        except Exception as e:
            print(f"Error generating embedding: {e}")
            return []

    def create_vector_id(self, file_path: str, section_title: str, chunk_index: int = 0) -> str:
        """Create unique vector ID in format: pageSlug#hN#cM."""
        # Extract page slug from file path or use file name
        page_slug = Path(file_path).stem
        if '/' in file_path and 'pages/' in file_path:
            page_slug = file_path.split('pages/')[-1].replace('.json', '')

        # Clean the section title to create a header identifier
        header_id = re.sub(r'[^a-zA-Z0-9]', '', section_title.lower())[:10]

        return f"{page_slug}#h{hash(section_title) % 1000}#c{chunk_index}"

    def process_documents(self, docs_dir: str = "docs") -> Dict[str, List[Dict[str, Any]]]:
        """Process all supported documents in the docs directory."""
        vectors_by_index = {}

        # Find all supported files
        supported_extensions = ['.md', '.json', '.txt']
        if PPTX_AVAILABLE:
            supported_extensions.extend(['.ppt', '.pptx'])
        if DOCX_AVAILABLE:
            supported_extensions.extend(['.doc', '.docx'])
        if PDF_AVAILABLE:
            supported_extensions.append('.pdf')

        all_files = []
        for ext in supported_extensions:
            all_files.extend(glob.glob(f"{docs_dir}/**/*{ext}", recursive=True))

        if not all_files:
            print(f"No supported files found in {docs_dir}")
            return vectors_by_index

        print(f"Found {len(all_files)} files: {len([f for f in all_files if f.endswith('.md')])} MD, {len([f for f in all_files if f.endswith('.json')])} JSON, {len([f for f in all_files if not f.endswith(('.md', '.json'))])} other")

        for file_path in all_files:
            print(f"\nProcessing: {file_path}")

            # Read file content with metadata
            content, file_metadata = self.read_file_content(file_path)
            if not content.strip():
                continue

            # Extract sections with enhanced metadata
            sections = self.extract_sections_from_content(content, file_path, file_metadata)
            print(f"  Extracted {len(sections)} sections")

            for section in sections:
                category = section["category"]

                # Only process if this category is in selected namespaces
                if category not in self.selected_namespaces:
                    print(f"    Skipping category '{category}' (not in selected namespaces)")
                    continue

                # Chunk the section content using token-based chunking
                chunks = self.chunk_text_by_tokens(section["content"])

                for chunk_index, chunk in enumerate(chunks):
                    if not chunk.strip():
                        continue

                    # Generate embedding
                    print(f"    Generating embedding for chunk {chunk_index + 1}/{len(chunks)} of '{section['title'][:50]}...' -> namespace: {category}")
                    embedding = self.generate_embedding(chunk)

                    if not embedding:
                        print(f"    Failed to generate embedding for chunk {chunk_index}")
                        continue

                    # Create vector
                    vector_id = self.create_vector_id(file_path, section["title"], chunk_index)

                    # Create enhanced metadata
                    metadata = {
                        "category": category,
                        "page_slug": section.get("page_slug", ""),
                        "heading_path": section.get("heading_path", section["title"]),
                        "source_file": section["source_file"],
                        "file_name": section["file_name"],
                        "source_url": section.get("source_url", ""),
                        "chunk_index": chunk_index,
                        "total_chunks": len(chunks),
                        "updated_at": datetime.now().isoformat(),
                        "text": chunk  # Store the actual text content
                    }

                    # Add category-specific metadata
                    if category == "ops":
                        metadata.update(self._extract_ops_metadata(chunk))
                    elif category == "medical":
                        metadata.update(self._extract_medical_metadata(chunk))
                    elif category == "legal":
                        metadata.update(self._extract_legal_metadata(chunk))
                    elif category == "program":
                        metadata.update(self._extract_program_metadata(chunk))

                    vector_data = {
                        "id": vector_id,
                        "values": embedding,
                        "metadata": metadata,
                        "namespace": category  # Store the target namespace
                    }

                    # Group by namespace instead of index
                    if category not in vectors_by_index:
                        vectors_by_index[category] = []
                    vectors_by_index[category].append(vector_data)

        return vectors_by_index

    def _extract_ops_metadata(self, content: str) -> Dict[str, Any]:
        """Extract operations-specific metadata from content."""
        metadata = {}
        content_lower = content.lower()

        # Extract dates
        date_patterns = [r'\b(\d{4}-\d{2}-\d{2})\b', r'\b(\w+ \d{1,2}, \d{4})\b']
        for pattern in date_patterns:
            dates = re.findall(pattern, content)
            if dates:
                metadata["dates_mentioned"] = dates[:3]  # Limit to first 3
                break

        # Extract pricing information
        price_patterns = [r'\$([0-9,]+)', r'([0-9,]+) dollars?']
        for pattern in price_patterns:
            prices = re.findall(pattern, content)
            if prices:
                try:
                    metadata["deposit_amount"] = int(prices[0].replace(',', ''))
                    break
                except ValueError:
                    pass

        # Extract location info
        if any(word in content_lower for word in ['address', 'location', 'colorado', 'denver', 'boulder']):
            metadata["has_location_info"] = True

        return metadata

    def _extract_medical_metadata(self, content: str) -> Dict[str, Any]:
        """Extract medical-specific metadata from content."""
        metadata = {"disclaimer_type": "medical", "requires_handoff": True}
        content_lower = content.lower()

        # Identify contraindications
        contraindications = []
        med_terms = ['ssri', 'maoi', 'antidepressant', 'blood thinner', 'heart condition', 'pregnancy', 'pregnant']
        for term in med_terms:
            if term in content_lower:
                contraindications.append(term)

        if contraindications:
            metadata["contraindications"] = contraindications

        # Check for screening requirement
        if any(word in content_lower for word in ['screening', 'intake', 'health assessment']):
            metadata["requires_screening"] = True

        return metadata

    def _extract_legal_metadata(self, content: str) -> Dict[str, Any]:
        """Extract legal-specific metadata from content."""
        metadata = {"disclaimer_type": "legal", "requires_handoff": False}
        content_lower = content.lower()

        # Check for Colorado-specific legal content
        if 'colorado' in content_lower:
            metadata["jurisdiction"] = "CO"

        # Check for licensing content
        if any(word in content_lower for word in ['license', 'licensed', 'rfra', 'religious freedom']):
            metadata["licensing_scope"] = "religious"

        # Check for privacy/data content
        if any(word in content_lower for word in ['privacy', 'data', 'personal information']):
            metadata["data_policy"] = True

        return metadata

    def _extract_program_metadata(self, content: str) -> Dict[str, Any]:
        """Extract program-specific metadata from content."""
        metadata = {}
        content_lower = content.lower()

        # Identify program type
        if 'retreat' in content_lower:
            metadata["program_type"] = "retreat"
        elif 'course' in content_lower:
            metadata["program_type"] = "course"
        elif 'training' in content_lower:
            metadata["program_type"] = "training"

        # Count ceremonies mentioned
        ceremony_count = content_lower.count('ceremony') + content_lower.count('ceremonies')
        if ceremony_count > 0:
            metadata["ceremonies_count"] = min(ceremony_count, 10)  # Cap at 10

        # Identify modalities
        modalities = []
        modality_terms = ['breathwork', 'yoga', 'meditation', 'integration', 'dance', 'sound healing']
        for term in modality_terms:
            if term in content_lower:
                modalities.append(term)

        if modalities:
            metadata["modalities"] = modalities

        return metadata

    def upsert_to_pinecone(self, vectors_by_namespace: Dict[str, List[Dict[str, Any]]], batch_size: int = 100):
        """Upsert vectors to Pinecone index across multiple namespaces."""
        total_vectors = sum(len(vectors) for vectors in vectors_by_namespace.values())

        if total_vectors == 0:
            print("No vectors to upsert")
            return

        print(f"\nUpserting {total_vectors} vectors to index '{self.index_name}'")
        print(f"Across {len([k for k, v in vectors_by_namespace.items() if v])} namespaces")

        try:
            index = self.connect_to_index()

            for namespace, vectors in vectors_by_namespace.items():
                if not vectors:
                    continue

                print(f"\n📁 Processing namespace '{namespace}': {len(vectors)} vectors")
                print(f"   📝 {self.selected_namespaces.get(namespace, 'Unknown category')}")

                # Upsert in batches
                for i in range(0, len(vectors), batch_size):
                    batch = vectors[i:i + batch_size]

                    try:
                        # Convert to the format Pinecone expects
                        pinecone_vectors = [
                            (vector["id"], vector["values"], vector["metadata"])
                            for vector in batch
                        ]

                        index.upsert(
                            vectors=pinecone_vectors,
                            namespace=namespace
                        )

                        print(f"  ✅ Upserted batch {i//batch_size + 1}/{(len(vectors)-1)//batch_size + 1} ({len(batch)} vectors)")

                    except Exception as e:
                        print(f"  ❌ Error upserting batch {i//batch_size + 1} to namespace {namespace}: {e}")

        except Exception as e:
            print(f"❌ Error connecting to index {self.index_name}: {e}")

    def verify_index_stats(self, specific_namespace: str = None):
        """Print statistics for the current index and namespaces."""
        if not self.index_name:
            print("❌ No index selected")
            return

        try:
            index = self.connect_to_index()
            stats = index.describe_index_stats()

            print(f"\n📊 Index: {self.index_name}")
            print(f"  Total vectors: {stats.get('total_vector_count', 0):,}")

            namespaces = stats.get('namespaces', {})

            if specific_namespace:
                if specific_namespace in namespaces:
                    ns_stats = namespaces[specific_namespace]
                    desc = self.selected_namespaces.get(specific_namespace, "Unknown")
                    print(f"\n📁 Namespace '{specific_namespace}' ({desc}):")
                    print(f"  Vectors: {ns_stats.get('vector_count', 0):,}")
                else:
                    print(f"  No vectors in '{specific_namespace}' namespace")
            else:
                if namespaces:
                    print(f"\n📁 All namespaces:")
                    for ns_name, ns_stats in namespaces.items():
                        desc = self.selected_namespaces.get(ns_name, "Unknown")
                        print(f"  • {ns_name}: {ns_stats.get('vector_count', 0):,} vectors ({desc})")
                else:
                    print("  No namespaces found")

        except Exception as e:
            print(f"❌ Error getting stats: {e}")

    def search_test(self, query: str = "What day is the cold plunge?", namespace: str = "ops", top_k: int = 3):
        """Test search functionality on a specific namespace."""
        print(f"\n🔍 Testing search on index '{self.index_name}', namespace '{namespace}'")
        print(f"Query: '{query}'")

        try:
            index = self.connect_to_index()

            # Generate embedding for query
            query_embedding = self.generate_embedding(query)
            if not query_embedding:
                print("Failed to generate query embedding")
                return

            # Search
            results = index.query(
                vector=query_embedding,
                namespace=namespace,
                top_k=top_k,
                include_metadata=True
            )

            print(f"\n📋 Found {len(results.matches)} results:")
            for i, match in enumerate(results.matches):
                metadata = match.metadata
                print(f"\n  {i+1}. 📊 Score: {match.score:.4f}")
                print(f"     📝 Page: {metadata.get('page_slug', 'N/A')}")
                print(f"     🏷️  Category: {metadata.get('category', 'N/A')}")
                print(f"     📍 Path: {metadata.get('heading_path', 'N/A')}")
                print(f"     💬 Content: {metadata.get('text', 'N/A')[:150]}...")

        except Exception as e:
            print(f"❌ Error during search test: {e}")

    def clear_namespace(self, index_name: str, namespace: str = "default", confirm: bool = False):
        """Clear vectors from a specific namespace in an index."""
        if not confirm:
            click.echo(f"This will delete ALL vectors in namespace '{namespace}' from index '{index_name}'")
            if not click.confirm("Are you sure?"):
                click.echo("Operation cancelled")
                return

        try:
            index = self.connect_to_index(index_name)

            # Get all vector IDs in the namespace
            stats = index.describe_index_stats()
            if namespace not in stats.get('namespaces', {}):
                click.echo(f"Namespace '{namespace}' not found in {index_name}")
                return

            # Delete the entire namespace
            index.delete(delete_all=True, namespace=namespace)
            click.echo(f"✅ Cleared namespace '{namespace}' from {index_name}")

        except Exception as e:
            click.echo(f"❌ Error clearing namespace: {e}")

# CLI Commands
@click.group()
@click.version_option(version="1.0.0")
def cli():
    """
    📚 Document Embedder CLI Tool

    Embed markdown documents using OpenAI text-embedding-3-small
    and store them in Pinecone for semantic search.
    """
    pass

@cli.command()
@click.argument('docs_path', default='docs')
@click.option('--index-name', help='Pinecone index name (will prompt if not provided)')
@click.option('--batch-size', default=100, help='Batch size for upserts')
@click.option('--max-tokens', default=500, help='Maximum tokens per chunk')
@click.option('--overlap-tokens', default=80, help='Token overlap between chunks')
@click.option('--auto-defaults', is_flag=True, help='Use default settings without prompting')
@click.option('--dry-run', is_flag=True, help='Show what would be processed without actually embedding')
def embed(docs_path, index_name, batch_size, max_tokens, overlap_tokens, auto_defaults, dry_run):
    """
    📄 Embed documents from DOCS_PATH into categorized Pinecone indices.

    Examples:
        embedder embed docs/
        embedder embed docs/website/ --namespace website
        embedder embed . --dry-run
    """
    if not _check_env_vars():
        return

    try:
        embedder = DocumentEmbedder()

        # Get index and namespace configuration
        if auto_defaults:
            embedder.index_name = index_name or "journeys-faq"
            embedder.selected_namespaces = embedder.default_namespaces.copy()
        else:
            if index_name:
                embedder.index_name = index_name
            else:
                embedder.index_name = embedder.prompt_for_index_selection()

            embedder.selected_namespaces = embedder.prompt_for_namespace_selection()

        print(f"\n✅ Configuration:")
        print(f"   📊 Index: {embedder.index_name}")
        print(f"   📁 Namespaces: {', '.join(embedder.selected_namespaces.keys())}")

        if dry_run:
            click.echo(f"🔍 DRY RUN - Scanning {docs_path} for supported files...")

            # Find all supported files
            supported_extensions = ['.md', '.json', '.txt']
            if PPTX_AVAILABLE:
                supported_extensions.extend(['.ppt', '.pptx'])
            if DOCX_AVAILABLE:
                supported_extensions.extend(['.doc', '.docx'])
            if PDF_AVAILABLE:
                supported_extensions.append('.pdf')

            all_files = []
            for ext in supported_extensions:
                all_files.extend(glob.glob(f"{docs_path}/**/*{ext}", recursive=True))

            if all_files:
                click.echo(f"📁 Would process {len(all_files)} files:")

                # Group by expected category for preview
                category_preview = {}
                for namespace in embedder.selected_namespaces.keys():
                    category_preview[namespace] = []

                for file_path in all_files:
                    category = embedder.categorize_content(file_path, "", "")
                    if category in category_preview:
                        category_preview[category].append(file_path)

                for category, files in category_preview.items():
                    if files:
                        desc = embedder.selected_namespaces[category]
                        click.echo(f"\n  📁 namespace '{category}' ({desc}) - {len(files)} files:")
                        for file_path in files[:5]:  # Show first 5
                            click.echo(f"    • {file_path}")
                        if len(files) > 5:
                            click.echo(f"    ... and {len(files) - 5} more")
            else:
                click.echo(f"❌ No supported files found in {docs_path}")
            return

        click.echo(f"🚀 Starting embedding process...")
        click.echo(f"   📁 Source: {docs_path}")
        click.echo(f"   📊 Index: {embedder.index_name}")
        click.echo(f"   🧩 Max tokens per chunk: {max_tokens}")
        click.echo(f"   🔗 Token overlap: {overlap_tokens}")

        # Show current stats
        click.echo(f"\n📊 Current index statistics:")
        embedder.verify_index_stats()

        # Process documents
        vectors_by_namespace = embedder.process_documents(docs_path)

        total_vectors = sum(len(vectors) for vectors in vectors_by_namespace.values())
        if total_vectors > 0:
            click.echo(f"\n📊 Generated {total_vectors} vectors across namespaces:")
            for namespace, vectors in vectors_by_namespace.items():
                if vectors:
                    desc = embedder.selected_namespaces.get(namespace, "Unknown")
                    click.echo(f"  • {namespace}: {len(vectors)} vectors ({desc})")

            # Upsert to Pinecone
            embedder.upsert_to_pinecone(vectors_by_namespace, batch_size)

            # Show final stats
            click.echo(f"\n📊 Final index statistics:")
            embedder.verify_index_stats()

            click.echo(f"\n✅ Successfully embedded {total_vectors} document chunks!")
        else:
            click.echo("❌ No vectors generated")

    except Exception as e:
        click.echo(f"❌ Error: {e}")

@cli.command()
@click.argument('query')
@click.option('--index-name', help='Pinecone index name (will prompt if not provided)')
@click.option('--namespace', help='Namespace to search in (will prompt if not provided)', type=click.Choice(['ops', 'legal', 'medical', 'program']))
@click.option('--top-k', default=5, help='Number of results to return')
@click.option('--score-threshold', default=0.0, type=float, help='Minimum similarity score')
def search(query, index_name, namespace, top_k, score_threshold):
    """
    🔍 Search for documents using semantic similarity.

    Examples:
        embedder search "What day is the cold plunge?" --namespace ops
        embedder search "medication contraindications" --namespace medical
        embedder search "ceremony schedule" --namespace program --top-k 3
    """
    if not _check_env_vars():
        return

    try:
        embedder = DocumentEmbedder()

        # Get index configuration
        if index_name:
            embedder.index_name = index_name
        else:
            embedder.index_name = embedder.prompt_for_index_selection()

        # Set default namespaces for reference
        embedder.selected_namespaces = embedder.default_namespaces.copy()

        # Get namespace to search
        if not namespace:
            print(f"\n📁 Available namespaces:")
            for key, desc in embedder.default_namespaces.items():
                print(f"  • {key}: {desc}")
            while True:
                namespace = input("\nSelect namespace to search: ").strip()
                if namespace in embedder.default_namespaces:
                    break
                print(f"Please enter one of: {', '.join(embedder.default_namespaces.keys())}")

        click.echo(f"🔍 Searching in index '{embedder.index_name}'")
        click.echo(f"   📝 Query: '{query}'")
        click.echo(f"   📁 Namespace: {namespace}")
        click.echo(f"   🎯 Description: {embedder.default_namespaces.get(namespace, 'Unknown')}")

        # Generate embedding for query
        query_embedding = embedder.generate_embedding(query)
        if not query_embedding:
            click.echo("❌ Failed to generate query embedding")
            return

        # Search
        index_connection = embedder.connect_to_index()
        results = index_connection.query(
            vector=query_embedding,
            namespace=namespace,
            top_k=top_k,
            include_metadata=True
        )

        # Filter by score threshold
        filtered_results = [
            match for match in results.matches
            if match.score >= score_threshold
        ]

        if not filtered_results:
            click.echo(f"❌ No results found above score threshold {score_threshold}")
            return

        click.echo(f"\n📋 Found {len(filtered_results)} results:")
        click.echo("-" * 80)

        for i, match in enumerate(filtered_results):
            metadata = match.metadata
            click.echo(f"\n🎯 Result {i+1}")
            click.echo(f"   📊 Score: {match.score:.4f}")
            click.echo(f"   📝 Page: {metadata.get('page_slug', 'N/A')}")
            click.echo(f"   🏷️  Category: {metadata.get('category', 'N/A')}")
            click.echo(f"   📍 Path: {metadata.get('heading_path', 'N/A')}")
            click.echo(f"   📄 File: {metadata.get('file_name', 'N/A')}")

            if metadata.get('source_url'):
                click.echo(f"   🔗 URL: {metadata['source_url']}")

            content = metadata.get('text', 'N/A')
            if len(content) > 250:
                content = content[:250] + "..."
            click.echo(f"   💬 Content: {content}")

    except Exception as e:
        click.echo(f"❌ Error during search: {e}")

@cli.command()
@click.option('--index-name', help='Pinecone index name (will prompt if not provided)')
@click.option('--namespace', help='Specific namespace to show stats for')
def stats(index_name, namespace):
    """
    📊 Show Pinecone index statistics.

    Examples:
        embedder stats
        embedder stats --index-name journeys-faq
        embedder stats --namespace ops
    """
    if not _check_env_vars():
        return

    try:
        embedder = DocumentEmbedder()
        embedder.selected_namespaces = embedder.default_namespaces.copy()

        # Get index configuration
        if index_name:
            embedder.index_name = index_name
        else:
            embedder.index_name = embedder.prompt_for_index_selection()

        # Show stats
        embedder.verify_index_stats(namespace)

    except Exception as e:
        click.echo(f"❌ Error getting stats: {e}")

@cli.command()
@click.option('--index-name', help='Pinecone index name (will prompt if not provided)')
@click.option('--namespace', required=True, help='Namespace to clear', type=click.Choice(['ops', 'legal', 'medical', 'program']))
@click.option('--yes', is_flag=True, help='Skip confirmation prompt')
def clear(index_name, namespace, yes):
    """
    🗑️  Clear vectors from a Pinecone namespace.

    Examples:
        embedder clear --namespace ops
        embedder clear --namespace medical --index-name journeys-faq
        embedder clear --namespace program --yes  # Skip confirmation
    """
    if not _check_env_vars():
        return

    try:
        embedder = DocumentEmbedder()
        embedder.selected_namespaces = embedder.default_namespaces.copy()

        # Get index configuration
        if index_name:
            embedder.index_name = index_name
        else:
            embedder.index_name = embedder.prompt_for_index_selection()

        embedder.clear_namespace(embedder.index_name, namespace, confirm=yes)

    except Exception as e:
        click.echo(f"❌ Error: {e}")

@cli.command()
def config():
    """
    ⚙️  Show current configuration and check environment.
    """
    click.echo("⚙️  Configuration Check:")
    click.echo("-" * 30)

    # Check environment variables
    openai_key = os.getenv('OPENAI_API_KEY')
    pinecone_key = os.getenv('PINECONE_API_KEY')

    click.echo(f"🔑 OPENAI_API_KEY: {'✅ Set' if openai_key else '❌ Not set'}")
    click.echo(f"🔑 PINECONE_API_KEY: {'✅ Set' if pinecone_key else '❌ Not set'}")

    if openai_key and pinecone_key:
        try:
            # Test OpenAI connection
            client = OpenAI(api_key=openai_key)
            models = client.models.list()
            click.echo(f"🤖 OpenAI: ✅ Connected")

            # Test Pinecone connection
            pc = Pinecone(api_key=pinecone_key)
            indexes = pc.list_indexes()
            click.echo(f"🌲 Pinecone: ✅ Connected")

            available_indexes = [idx.name for idx in indexes]
            click.echo(f"📋 Available indexes: {available_indexes}")

            # Check if journeys-faq index exists
            if "journeys-faq" in available_indexes:
                click.echo(f"✅ Default index 'journeys-faq' found")
            else:
                click.echo(f"⚠️  Default index 'journeys-faq' not found")
                click.echo("💡 You can select a different index during embedding or create 'journeys-faq' in Pinecone console")

        except Exception as e:
            click.echo(f"❌ Connection error: {e}")
    else:
        click.echo("\n💡 To set up:")
        click.echo("   1. Copy .env.example to .env")
        click.echo("   2. Add your API keys to .env")
        click.echo("   3. Create a Pinecone index (e.g., 'journeys-faq')")
        click.echo("   4. Run: embedder config")

def _check_env_vars() -> bool:
    """Check if required environment variables are set."""
    openai_key = os.getenv('OPENAI_API_KEY')
    pinecone_key = os.getenv('PINECONE_API_KEY')

    if not openai_key:
        click.echo("❌ OPENAI_API_KEY environment variable not set")
        click.echo("💡 Set it in .env file or export OPENAI_API_KEY=your_key")
        return False

    if not pinecone_key:
        click.echo("❌ PINECONE_API_KEY environment variable not set")
        click.echo("💡 Set it in .env file or export PINECONE_API_KEY=your_key")
        return False

    return True

if __name__ == "__main__":
    cli()